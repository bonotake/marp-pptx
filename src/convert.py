from dotenv import load_dotenv
import os
from pptx import Presentation
from pptx.util import Inches, Pt
import re
from pathlib import Path
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def convert_marp_to_pptx(input_file: Path, output_file: Path, debug: bool = False) -> None:
    # デバッグ出力用のヘルパー関数
    def debug_print(*args, **kwargs):
        if debug:
            print(*args, **kwargs)
    
    # 設定値
    INDENT_SPACES = 2  # インデントの字数（Marpのデフォルトは2スペース）
    
    def apply_text_styles(p, text):
        """段落にテキストスタイルを適用する"""
        # スタイルマーカーとその適用方法を定義
        style_markers = [
            ('**', lambda run: setattr(run.font, 'bold', True)),         # 太字
            ('~~', lambda run: setattr(run.font, 'strike', True)),       # 打消し線
            ('*', lambda run: setattr(run.font, 'italic', True))         # イタリック
        ]
        
        # 現在のテキストとその位置を追跡
        current_text = text
        current_pos = 0
        p.text = ""  # 段落を空にする
        
        while current_text:
            # 最も近いマーカーとその位置を見つける
            next_marker = None
            next_pos = len(current_text)
            
            for marker, _ in style_markers:
                pos = current_text.find(marker)
                if pos != -1 and pos < next_pos:
                    next_marker = marker
                    next_pos = pos
            
            # マーカーが見つからない場合、残りのテキストを追加して終了
            if next_marker is None:
                if current_text:
                    run = p.add_run()
                    run.text = current_text
                break
            
            # マーカーまでのテキストを追加
            if next_pos > 0:
                run = p.add_run()
                run.text = current_text[:next_pos]
            
            # マーカーの終わりを探す
            end_pos = current_text.find(next_marker, next_pos + len(next_marker))
            if end_pos == -1:  # 閉じマーカーが見つからない
                run = p.add_run()
                run.text = current_text[next_pos:]
                break
            
            # スタイル適用されたテキストを追加
            styled_text = current_text[next_pos + len(next_marker):end_pos]
            run = p.add_run()
            run.text = styled_text
            
            # スタイルを適用
            for marker, apply_style in style_markers:
                if marker == next_marker:
                    apply_style(run)
            
            # 残りのテキストを更新
            current_text = current_text[end_pos + len(next_marker):]
    
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Marpファイルを読み込む
    with open(input_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # スライドを分割
    slides = content.split('---')
    
    # フロントマターを除去（最初の要素を無視）
    slides = slides[2:]  # 最初の2つの要素（空要素とフロントマター）を除外
    
    # CSSスタイルを解析
    style_definitions = {}
    css_pattern = re.compile(r'<style[^>]*>(.*?)</style>', re.DOTALL)
    css_matches = css_pattern.findall(content)
    
    for css in css_matches:
        # クラス定義を探す（例: .class-name { property: value; }）
        class_pattern = re.compile(r'\.([^{]+){([^}]+)}')
        for match in class_pattern.finditer(css):
            class_name = match.group(1).strip()
            properties = match.group(2).strip()
            style_definitions[class_name] = dict(
                prop.strip().split(':') 
                for prop in properties.split(';') 
                if ':' in prop
            )
    
    for i, slide_content in enumerate(slides):
        # 空のスライドをスキップ
        if not slide_content.strip():
            continue
        
        debug_print(f"=== Slide {i} ===")
        
        # 全体を1つの文字列として処理
        content = '\n'.join(slide_content.strip().split('\n'))
        div_pattern = re.compile(r'<div\s+(?:class|style)=["\']([^"\']+)["\']>(.*?)</div>', re.DOTALL)
        
        # divタグを処理
        while True:
            div_match = div_pattern.search(content)
            if not div_match:
                break
            
            class_name = div_match.group(1)
            div_content = div_match.group(2).strip()
            # divタグを内容で置換
            content = content[:div_match.start()] + div_content + content[div_match.end():]
        
        # 処理後の内容を行に分割
        lines = content.split('\n')
        
        # レイアウトの選択
        non_empty_lines = [line for line in lines if line.strip()]
        if i == 0:
            layout = prs.slide_layouts[0]  # タイトルスライド
        elif len(non_empty_lines) == 1 and non_empty_lines[0].strip().startswith('# '):
            layout = prs.slide_layouts[2]  # セクションヘッダー
        else:
            layout = prs.slide_layouts[1]  # タイトルと本文
            
        slide = prs.slides.add_slide(layout)
        
        # タイトルを探す（# で始まる最初の行）
        title = ""
        content_lines = []
        
        # タイトルプレースホルダーを取得
        title_shape = slide.shapes.title
        text_frame = None
        
        for line in lines:
            # h1（#）をスライドタイトルとして処理
            if not title and line.strip().startswith('# '):
                debug_print(f"Found title: '{line}'")
                title = line.strip('#').strip()
            # h2-h6（## ～ ######）を本文の見出しとして処理
            elif line.strip().startswith('#'):
                debug_print(f"Found heading: '{line}'")
                heading_level = len(line.split()[0])
                content = line.strip('#').strip()
                
                # 本文用のテキストフレームを初期化（まだ作られていない場合）
                if not text_frame and slide.placeholders[1]:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.word_wrap = True
                
                if text_frame:
                    p = text_frame.add_paragraph()
                    apply_text_styles(p, content)
                    # 見出しレベルに応じてスタイルを設定
                    if heading_level == 2:  # h2
                        debug_print(f"  Setting h2 style for: '{content}'")
                        p.font.size = Pt(32)
                        p.font.bold = True
                    elif heading_level == 3:  # h3
                        debug_print(f"  Setting h3 style for: '{content}'")
                        p.font.size = Pt(28)
                        p.font.bold = True
                    elif heading_level <= 6:  # h4-h6
                        debug_print(f"  Setting h4-6 style for: '{content}'")
                        p.font.size = Pt(24)
                        p.font.bold = True
                else:
                    debug_print("Warning: text_frame is not initialized")
            else:
                debug_print(f"Adding to content: '{line}'")
                content_lines.append(line)
        
        # タイトルを設定
        if title and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.left = Inches(0.5)
            title_shape.top = Inches(0.5)
            title_shape.width = Inches(15)  # 16 - (0.5 * 2)
            title_shape.text = title
            
        # 本文を設定
        if content_lines and slide.placeholders[1]:
            body_shape = slide.placeholders[1]
            body_shape.left = Inches(0.5)
            body_shape.top = Inches(2)
            body_shape.width = Inches(15)  # 16 - (0.5 * 2)
            body_shape.height = Inches(6.5)  # 9 - 2 - 0.5
            
            text_frame = body_shape.text_frame
            text_frame.word_wrap = True
            
            for line in content_lines:
                p = text_frame.add_paragraph()
                if line.strip().startswith('- '):
                    indent_level = (len(line) - len(line.lstrip())) // INDENT_SPACES
                    text = line.strip('- ').strip()
                    apply_text_styles(p, text)
                    p.level = indent_level
                    p.bullet = True
                else:
                    text = line.strip()
                    apply_text_styles(p, text)
    
    # PowerPointファイルを保存
    prs.save(str(output_file))

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Convert Marp markdown to PowerPoint')
    parser.add_argument('--debug', action='store_true', help='Enable debug output')
    args = parser.parse_args()
    
    load_dotenv()
    work_folder: Path = Path(os.getenv("WORK_FOLDER"))
    input_file: Path = work_folder / Path("main.md")
    output_file: Path = work_folder / Path("presentation.pptx")
    convert_marp_to_pptx(input_file, output_file, debug=args.debug)

